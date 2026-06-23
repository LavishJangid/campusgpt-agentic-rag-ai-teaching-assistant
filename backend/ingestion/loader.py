"""Document loader module - extracts text from PDF, DOCX, PPTX, and TXT files."""

import os
import re
from pathlib import Path
from typing import Optional

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from loguru import logger


class DocumentLoader:
    """Handles text extraction from multiple document formats."""

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".pptx"}

    def __init__(self):
        self.stats = {"processed": 0, "failed": 0, "total_pages": 0}

    def load(self, file_path: str) -> dict:
        """
        Load a document and extract text with metadata.

        Returns:
            dict with keys: text, metadata, pages
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported format: {ext}. Supported: {self.SUPPORTED_EXTENSIONS}")

        logger.info(f"Loading document: {path.name} ({ext})")

        try:
            if ext == ".pdf":
                result = self._load_pdf(path)
            elif ext == ".docx":
                result = self._load_docx(path)
            elif ext == ".pptx":
                result = self._load_pptx(path)
            elif ext == ".txt":
                result = self._load_txt(path)
            else:
                raise ValueError(f"Unsupported format: {ext}")

            self.stats["processed"] += 1
            logger.info(f"Successfully loaded: {path.name} | Pages: {len(result['pages'])}")
            return result

        except Exception as e:
            self.stats["failed"] += 1
            logger.error(f"Failed to load {path.name}: {str(e)}")
            raise

    def _load_pdf(self, path: Path) -> dict:
        """Extract text from PDF with page-level granularity."""
        reader = PdfReader(str(path))
        pages = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = self._clean_text(text)
            if text.strip():
                pages.append({
                    "text": text,
                    "page_number": i + 1,
                    "source": path.name,
                })

        self.stats["total_pages"] += len(pages)
        full_text = "\n\n".join([p["text"] for p in pages])

        return {
            "text": full_text,
            "pages": pages,
            "metadata": {
                "source": path.name,
                "file_type": "pdf",
                "total_pages": len(reader.pages),
                "extracted_pages": len(pages),
                "file_size_kb": round(path.stat().st_size / 1024, 2),
            },
        }

    def _load_docx(self, path: Path) -> dict:
        """Extract text from DOCX documents."""
        doc = DocxDocument(str(path))
        pages = []
        current_text = []
        page_num = 1

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                current_text.append(text)

            # Simulate page breaks every ~3000 chars
            joined = "\n".join(current_text)
            if len(joined) > 3000:
                pages.append({
                    "text": self._clean_text(joined),
                    "page_number": page_num,
                    "source": path.name,
                })
                current_text = []
                page_num += 1

        # Add remaining text
        if current_text:
            pages.append({
                "text": self._clean_text("\n".join(current_text)),
                "page_number": page_num,
                "source": path.name,
            })

        self.stats["total_pages"] += len(pages)
        full_text = "\n\n".join([p["text"] for p in pages])

        return {
            "text": full_text,
            "pages": pages,
            "metadata": {
                "source": path.name,
                "file_type": "docx",
                "total_pages": len(pages),
                "extracted_pages": len(pages),
                "file_size_kb": round(path.stat().st_size / 1024, 2),
            },
        }

    def _load_pptx(self, path: Path) -> dict:
        """Extract text from PowerPoint presentations."""
        prs = Presentation(str(path))
        pages = []

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())

                # Also extract from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            texts.append(row_text)

            slide_text = "\n".join(texts)
            if slide_text.strip():
                pages.append({
                    "text": self._clean_text(slide_text),
                    "page_number": i + 1,
                    "source": path.name,
                })

        self.stats["total_pages"] += len(pages)
        full_text = "\n\n".join([p["text"] for p in pages])

        return {
            "text": full_text,
            "pages": pages,
            "metadata": {
                "source": path.name,
                "file_type": "pptx",
                "total_pages": len(prs.slides),
                "extracted_pages": len(pages),
                "file_size_kb": round(path.stat().st_size / 1024, 2),
            },
        }

    def _load_txt(self, path: Path) -> dict:
        """Extract text from plain text files."""
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
        except UnicodeDecodeError:
            import chardet
            with open(path, "rb") as f:
                raw = f.read()
            detected = chardet.detect(raw)
            content = raw.decode(detected["encoding"] or "utf-8", errors="replace")

        content = self._clean_text(content)
        lines = content.split("\n")

        # Split into pages of ~50 lines
        page_size = 50
        pages = []
        for i in range(0, len(lines), page_size):
            chunk = "\n".join(lines[i : i + page_size])
            if chunk.strip():
                pages.append({
                    "text": chunk,
                    "page_number": (i // page_size) + 1,
                    "source": path.name,
                })

        self.stats["total_pages"] += len(pages)

        return {
            "text": content,
            "pages": pages,
            "metadata": {
                "source": path.name,
                "file_type": "txt",
                "total_pages": len(pages),
                "extracted_pages": len(pages),
                "file_size_kb": round(path.stat().st_size / 1024, 2),
            },
        }

    @staticmethod
    def _clean_text(text: str) -> str:
        """Clean and normalize extracted text."""
        # Remove excessive whitespace
        text = re.sub(r"\s+", " ", text)
        # Remove special characters but keep punctuation
        text = re.sub(r"[^\w\s.,;:!?()'\"-/\\@#$%&*+=<>\[\]{}|~`^]", "", text)
        # Fix double spaces
        text = re.sub(r"  +", " ", text)
        # Fix line breaks
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()
